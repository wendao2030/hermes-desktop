#!/usr/bin/env python3
"""
PowerPoint 生成器 — v2.0.0

纯 python-pptx 实现，生成标准可编辑 PPTX。
7 种幻灯片布局、5 套配色方案、图表、表格、图片、渐变背景、页码页脚。

用法:
  python3 generate_pptx.py <json_input> [--out output.pptx] [--style business_blue]
  python3 generate_pptx.py --interactive
"""

import os
import sys
import json

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData


# ── Color palettes ──────────────────────────────────────
PALETTES = {
    "business_blue": {
        "primary": RGBColor(30, 60, 114),
        "secondary": RGBColor(70, 130, 180),
        "accent": RGBColor(255, 193, 7),
        "highlight": RGBColor(231, 76, 60),
        "success": RGBColor(39, 174, 96),
        "warning": RGBColor(230, 126, 34),
        "text_primary": RGBColor(44, 62, 80),
        "text_secondary": RGBColor(127, 140, 141),
        "bg_primary": RGBColor(255, 255, 255),
        "bg_secondary": RGBColor(236, 240, 241),
        "gradient_start": RGBColor(30, 60, 114),
        "gradient_end": RGBColor(52, 152, 219),
    },
    "academic_white": {
        "primary": RGBColor(0, 51, 102),
        "secondary": RGBColor(100, 100, 100),
        "accent": RGBColor(178, 34, 34),
        "highlight": RGBColor(0, 102, 153),
        "success": RGBColor(34, 139, 34),
        "warning": RGBColor(184, 134, 11),
        "text_primary": RGBColor(51, 51, 51),
        "text_secondary": RGBColor(119, 119, 119),
        "bg_primary": RGBColor(255, 255, 255),
        "bg_secondary": RGBColor(245, 245, 245),
        "gradient_start": RGBColor(0, 51, 102),
        "gradient_end": RGBColor(0, 102, 153),
    },
    "creative_purple": {
        "primary": RGBColor(102, 45, 140),
        "secondary": RGBColor(155, 89, 182),
        "accent": RGBColor(241, 196, 15),
        "highlight": RGBColor(231, 76, 60),
        "success": RGBColor(46, 204, 113),
        "warning": RGBColor(230, 126, 34),
        "text_primary": RGBColor(52, 31, 71),
        "text_secondary": RGBColor(136, 114, 158),
        "bg_primary": RGBColor(248, 244, 252),
        "bg_secondary": RGBColor(235, 225, 245),
        "gradient_start": RGBColor(102, 45, 140),
        "gradient_end": RGBColor(155, 89, 182),
    },
    "tech_dark": {
        "primary": RGBColor(48, 56, 72),
        "secondary": RGBColor(88, 100, 120),
        "accent": RGBColor(0, 200, 150),
        "highlight": RGBColor(0, 180, 216),
        "success": RGBColor(46, 204, 113),
        "warning": RGBColor(241, 196, 15),
        "text_primary": RGBColor(230, 235, 240),
        "text_secondary": RGBColor(170, 180, 195),
        "bg_primary": RGBColor(60, 68, 85),
        "bg_secondary": RGBColor(75, 85, 105),
        "gradient_start": RGBColor(48, 56, 72),
        "gradient_end": RGBColor(70, 80, 100),
    },
    "minimal_gray": {
        "primary": RGBColor(80, 80, 80),
        "secondary": RGBColor(150, 150, 150),
        "accent": RGBColor(0, 120, 215),
        "highlight": RGBColor(215, 80, 60),
        "success": RGBColor(60, 170, 100),
        "warning": RGBColor(210, 160, 40),
        "text_primary": RGBColor(51, 51, 51),
        "text_secondary": RGBColor(153, 153, 153),
        "bg_primary": RGBColor(250, 250, 250),
        "bg_secondary": RGBColor(240, 240, 240),
        "gradient_start": RGBColor(80, 80, 80),
        "gradient_end": RGBColor(140, 140, 140),
    },
}


def _c(name, pal):
    return pal.get(name, PALETTES["business_blue"][name])


# ── Helpers ──────────────────────────────────────────────
def _add_textbox(slide, left, top, w, h):
    return slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))


def _add_para(tf, text, size=18, color=None, bold=False, align=None, spacing=0, level=0):
    if tf.paragraphs and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.level = level
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    if spacing:
        p.space_after = Pt(spacing)
    return p


def _add_shape(slide, shape_type, left, top, w, h, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    return shape


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_page_number(slide, num, pal, prs):
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Inches(0.3), Emu(sw) - Inches(0.4), Inches(1), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(8)
    p.font.color.rgb = _c("text_secondary", pal)
    p.alignment = PP_ALIGN.LEFT


def _add_footer(slide, text, pal, prs):
    if not text:
        return
    sw = prs.slide_width
    box = slide.shapes.add_textbox(Inches(5), Emu(sw) - Inches(0.4), Inches(5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = _c("text_secondary", pal)
    p.alignment = PP_ALIGN.RIGHT


def _validate_file(path):
    if not os.path.exists(path):
        raise FileNotFoundError(f"文件不存在: {path}")


# ── Generator ────────────────────────────────────────────
class PPTGenerator:
    SLIDE_W = 13.333
    SLIDE_H = 7.5

    def __init__(self, style="business_blue", footer="", page_numbers=True):
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.SLIDE_W)
        self.prs.slide_height = Inches(self.SLIDE_H)
        self.pal = PALETTES.get(style, PALETTES["business_blue"])
        self.style = style
        self.footer = footer
        self.page_numbers = page_numbers
        self._page = 0

    def _next_page(self):
        self._page += 1
        return self._page

    def _make_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        pn = self._next_page()
        if self.page_numbers:
            _add_page_number(slide, pn, self.pal, self.prs)
        if self.footer:
            _add_footer(slide, self.footer, self.pal, self.prs)
        return slide

    def _title_bar(self, slide, title, height=1):
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_W, height,
                   fill_color=_c("primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, height, self.SLIDE_W, 0.06,
                   fill_color=_c("accent", self.pal))
        tb = _add_textbox(slide, 0.8, 0.12, 11.733, height - 0.2)
        _add_para(tb.text_frame, title, size=26, bold=True,
                  color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

    def add_cover(self, title, subtitle="", author="", date="", variant="centered"):
        slide = self._make_slide()
        if variant == "split":
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 6.667, self.SLIDE_H,
                       fill_color=_c("primary", self.pal))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 6.667, 0, 0.06, self.SLIDE_H,
                       fill_color=_c("accent", self.pal))
            tb = _add_textbox(slide, 7.2, 2, 5.3, 4)
            _add_para(tb.text_frame, title, size=38, bold=True, color=_c("bg_primary", self.pal),
                      align=PP_ALIGN.LEFT, spacing=14)
            if subtitle:
                _add_para(tb.text_frame, subtitle, size=18, color=_c("bg_secondary", self.pal),
                          align=PP_ALIGN.LEFT)
        elif variant == "left":
            _set_bg(slide, _c("bg_primary", self.pal))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.6, self.SLIDE_H,
                       fill_color=_c("primary", self.pal))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.6, 0, 0.06, self.SLIDE_H,
                       fill_color=_c("accent", self.pal))
            tb = _add_textbox(slide, 1.5, 2.2, 10.8, 3)
            _add_para(tb.text_frame, title, size=40, bold=True, color=_c("primary", self.pal),
                      align=PP_ALIGN.LEFT, spacing=12)
            if subtitle:
                _add_para(tb.text_frame, subtitle, size=20, color=_c("text_secondary", self.pal),
                          align=PP_ALIGN.LEFT)
        else:
            _set_bg(slide, _c("bg_primary", self.pal))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, 3.5, 6.333, 0.06,
                       fill_color=_c("accent", self.pal))
            tb = _add_textbox(slide, 1.5, 1.8, 10.333, 2.5)
            _add_para(tb.text_frame, title, size=42, bold=True, color=_c("primary", self.pal),
                      align=PP_ALIGN.CENTER, spacing=14)
            if subtitle:
                _add_para(tb.text_frame, subtitle, size=20, color=_c("secondary", self.pal),
                          align=PP_ALIGN.CENTER)

        if author or date:
            meta = "  |  ".join(filter(None, [author, date]))
            y = 6.2 if variant != "split" else 6
            mb = _add_textbox(slide, 1.5, y, 10.333, 0.5)
            _add_para(mb.text_frame, meta, size=13,
                      color=_c("text_secondary", self.pal), align=PP_ALIGN.CENTER)
        return slide

    def add_section(self, title, subtitle=""):
        slide = self._make_slide()
        _set_bg(slide, _c("primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 4, 3.4, 5.333, 0.06,
                   fill_color=_c("accent", self.pal))
        tb = _add_textbox(slide, 1.5, 1.2, 10.333, 2)
        _add_para(tb.text_frame, title, size=34, bold=True, color=_c("bg_primary", self.pal),
                  align=PP_ALIGN.CENTER, spacing=10)
        if subtitle:
            _add_para(tb.text_frame, subtitle, size=17, color=_c("bg_secondary", self.pal),
                      align=PP_ALIGN.CENTER)
        return slide

    def add_content(self, title, items, columns=1):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        if columns == 1:
            tb = _add_textbox(slide, 0.9, 1.5, 11.533, 5.3)
            for item in items:
                if isinstance(item, dict):
                    _add_para(tb.text_frame, item.get("t", ""), size=20, bold=True,
                              color=_c("primary", self.pal), spacing=4)
                    for sub in item.get("s", []):
                        _add_para(tb.text_frame, f"  {sub}", size=15,
                                  color=_c("text_secondary", self.pal), spacing=3, level=1)
                else:
                    _add_para(tb.text_frame, str(item), size=17,
                              color=_c("text_primary", self.pal), spacing=9)
        elif columns == 2:
            col_w = 5.6
            sep_x = 6.6
            _add_shape(slide, MSO_SHAPE.RECTANGLE, sep_x, 1.3, 0.03, 5.5,
                       fill_color=_c("bg_secondary", self.pal))
            for ci in range(2):
                tb = _add_textbox(slide, 0.9 + ci * (col_w + 1), 1.4, col_w, 5.3)
                for item in items[ci::2]:
                    _add_para(tb.text_frame, str(item), size=15,
                              color=_c("text_primary", self.pal), spacing=7)
        return slide

    def add_two_column(self, title, left_items, right_items):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 6.6, 1.3, 0.03, 5.5,
                   fill_color=_c("bg_secondary", self.pal))
        for ci, (items) in enumerate([left_items, right_items]):
            tb = _add_textbox(slide, 0.9 + ci * 6, 1.4, 5.3, 5.3)
            for item in items:
                _add_para(tb.text_frame, str(item), size=15,
                          color=_c("text_primary", self.pal), spacing=7)
        return slide

    def add_table(self, title, headers, rows):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        nrows = len(rows) + 1
        ncols = len(headers)
        if ncols == 0 or nrows == 0:
            return slide
        tbl = slide.shapes.add_table(nrows, ncols, Inches(0.5), Inches(1.5),
                                     Inches(12.333), Inches(5.3)).table
        for i, h in enumerate(headers):
            c = tbl.cell(0, i)
            c.text = str(h)
            c.fill.solid()
            c.fill.fore_color.rgb = _c("primary", self.pal)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                p.alignment = PP_ALIGN.CENTER
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                c = tbl.cell(ri, ci)
                c.text = str(val)
                if ri % 2 == 0:
                    c.fill.solid()
                    c.fill.fore_color.rgb = _c("bg_secondary", self.pal)
                for p in c.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.color.rgb = _c("text_primary", self.pal)
                    p.alignment = PP_ALIGN.CENTER
        return slide

    def add_image_slide(self, title, image_path, caption="", overlay=False):
        _validate_file(image_path)
        slide = self._make_slide()
        if overlay:
            slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                                     Inches(self.SLIDE_W), Inches(self.SLIDE_H))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_W, self.SLIDE_H,
                       fill_color=_c("primary", self.pal))
            _add_shape(slide, MSO_SHAPE.RECTANGLE, 3.5, 4.8, 6.333, 0.06,
                       fill_color=_c("accent", self.pal))
            tb = _add_textbox(slide, 1.5, 2, 10.333, 2.5)
            _add_para(tb.text_frame, title, size=38, bold=True,
                      color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, spacing=10)
            if caption:
                _add_para(tb.text_frame, caption, size=18,
                          color=_c("bg_secondary", self.pal), align=PP_ALIGN.CENTER)
        else:
            _set_bg(slide, _c("bg_primary", self.pal))
            self._title_bar(slide, title)
            slide.shapes.add_picture(image_path, Inches(1.2), Inches(1.5),
                                     Inches(10.933), Inches(4.8))
            if caption:
                tb = _add_textbox(slide, 0.9, 6.4, 11.533, 0.5)
                _add_para(tb.text_frame, caption, size=12,
                          color=_c("text_secondary", self.pal), align=PP_ALIGN.CENTER)
        return slide

    def add_summary(self, title, points, conclusion=""):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        tb = _add_textbox(slide, 1.2, 1.5, 10.933, 4)
        for pt in points:
            txt = pt if isinstance(pt, str) else pt.get("t", str(pt))
            _add_para(tb.text_frame, txt, size=17,
                      color=_c("text_primary", self.pal), spacing=10)
        if conclusion:
            box = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1, 5.6, 11.333, 1,
                             fill_color=_c("secondary", self.pal))
            tf = box.text_frame
            tf.word_wrap = True
            _add_para(tf, conclusion, size=17, bold=True,
                      color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)
        return slide

    def add_timeline(self, title, milestones):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        n = len(milestones)
        if n == 0:
            return slide
        spacing = 11.333 / max(n, 1)
        line_y = 3.8
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0.8, line_y, 11.733, 0.06,
                   fill_color=_c("secondary", self.pal))
        for i, ms in enumerate(milestones):
            cx = 0.8 + spacing * i + spacing / 2 - 0.2
            if isinstance(ms, dict):
                label = ms.get("label", "")
                desc = ms.get("desc", "")
                active = ms.get("active", i == 0)
            else:
                label = str(ms)
                desc = ""
                active = False
            dot_color = _c("accent", self.pal) if active else _c("secondary", self.pal)
            _add_shape(slide, MSO_SHAPE.OVAL, cx, line_y - 0.2, 0.4, 0.4,
                       fill_color=dot_color)
            lb = _add_textbox(slide, cx - 0.8, line_y - 1.3, 2, 0.9)
            _add_para(lb.text_frame, label, size=13, bold=True,
                      color=_c("primary", self.pal), align=PP_ALIGN.CENTER, spacing=2)
            if desc:
                db = _add_textbox(slide, cx - 0.8, line_y + 0.7, 2, 0.9)
                _add_para(db.text_frame, desc, size=11,
                          color=_c("text_secondary", self.pal), align=PP_ALIGN.CENTER)
        return slide

    def add_quote(self, quote_text, attribution=""):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        qbox = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 1.5, 10.933, 4.5,
                          fill_color=_c("bg_secondary", self.pal))
        qbox.line.color.rgb = _c("secondary", self.pal)
        qbox.line.width = Pt(1.5)
        ib = _add_textbox(slide, 1.2, 1.3, 1, 1)
        _add_para(ib.text_frame, "❝", size=44, color=_c("secondary", self.pal),
                  align=PP_ALIGN.CENTER)
        tb = _add_textbox(slide, 2.2, 2, 8.933, 3)
        _add_para(tb.text_frame, quote_text, size=22, color=_c("text_primary", self.pal),
                  align=PP_ALIGN.LEFT, spacing=8)
        if attribution:
            ab = _add_textbox(slide, 2.2, 4.8, 8.933, 0.7)
            _add_para(ab.text_frame, f"— {attribution}", size=15,
                      color=_c("text_secondary", self.pal), align=PP_ALIGN.RIGHT)
        return slide

    def add_contact(self, title, contact_info):
        slide = self._make_slide()
        _set_bg(slide, _c("primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 4, 3.4, 5.333, 0.06,
                   fill_color=_c("accent", self.pal))
        tb = _add_textbox(slide, 1.5, 1.5, 10.333, 1.5)
        _add_para(tb.text_frame, title, size=34, bold=True,
                  color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, spacing=14)
        if isinstance(contact_info, list):
            contact_info = "\n".join(contact_info)
        cb = _add_textbox(slide, 2, 3.8, 9.333, 3)
        _add_para(cb.text_frame, contact_info, size=17,
                  color=_c("bg_secondary", self.pal), align=PP_ALIGN.CENTER, spacing=7)
        return slide

    def add_agenda(self, title, topics, highlight=0):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        tb = _add_textbox(slide, 2.5, 1.8, 8.333, 5)
        for i, t in enumerate(topics):
            num = f"{i + 1:02d}"
            is_hl = i == highlight
            _add_para(tb.text_frame, f"{num}    {t}", size=19,
                      bold=is_hl, color=_c("accent", self.pal) if is_hl else _c("text_primary", self.pal),
                      align=PP_ALIGN.LEFT, spacing=13)
        return slide

    def add_chart(self, title, chart_type, categories, series_data):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        self._title_bar(slide, title)
        cd = CategoryChartData()
        cd.categories = categories
        for s in series_data:
            cd.add_series(s.get("name", ""), s.get("values", []))
        ct_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        ct = ct_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_shape = slide.shapes.add_chart(
            ct, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.3), cd
        )
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        if hasattr(chart, 'font'):
            chart.font.size = Pt(10)
        if chart_type == "pie":
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.font.size = Pt(10)
        return slide


    def add_quote(self, quote_text, attribution=""):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        qbox = _add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.2, 1.5, 10.933, 4.5,
                          fill_color=_c("bg_secondary", self.pal))
        qbox.line.color.rgb = _c("secondary", self.pal)
        qbox.line.width = Pt(1.5)
        ib = _add_textbox(slide, 1.2, 1.3, 1, 1)
        _add_para(ib.text_frame, "❝", size=44, color=_c("secondary", self.pal),
                  align=PP_ALIGN.CENTER)
        tb = _add_textbox(slide, 2.2, 2, 8.933, 3)
        _add_para(tb.text_frame, quote_text, size=22, color=_c("text_primary", self.pal),
                  align=PP_ALIGN.LEFT, spacing=8)
        if attribution:
            ab = _add_textbox(slide, 2.2, 4.8, 8.933, 0.7)
            _add_para(ab.text_frame, f"— {attribution}", size=15,
                      color=_c("text_secondary", self.pal), align=PP_ALIGN.RIGHT)
        return slide

    def add_contact(self, title, contact_info):
        slide = self._make_slide()
        _set_bg(slide, _c("primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 4, 3.4, 5.333, 0.06,
                   fill_color=_c("accent", self.pal))
        tb = _add_textbox(slide, 1.5, 1.5, 10.333, 1.5)
        _add_para(tb.text_frame, title, size=34, bold=True,
                  color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER, spacing=14)
        if isinstance(contact_info, list):
            contact_info = "\n".join(contact_info)
        cb = _add_textbox(slide, 2, 3.8, 9.333, 3)
        _add_para(cb.text_frame, contact_info, size=17,
                  color=_c("bg_secondary", self.pal), align=PP_ALIGN.CENTER, spacing=7)
        return slide

    def add_agenda(self, title, topics, highlight=0):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_W, 1,
                   fill_color=_c("primary", self.pal))
        tb = _add_textbox(slide, 0.8, 0.12, 11.733, 0.8)
        _add_para(tb.text_frame, title, size=28, bold=True,
                  color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

        tb = _add_textbox(slide, 1.5, 1.8, 10.333, 5)
        for i, t in enumerate(topics):
            num = f"{i + 1:02d}"
            prefix = f"{num}  "
            if i == highlight:
                c = _c("accent", self.pal)
                b = True
            else:
                c = _c("text_primary", self.pal)
                b = False
            _add_para(tb.text_frame, f"{prefix}{t}", size=20, bold=b,
                      color=c, align=PP_ALIGN.LEFT, spacing=14)
        return slide

    def add_chart(self, title, chart_type, categories, series_data):
        slide = self._make_slide()
        _set_bg(slide, _c("bg_primary", self.pal))
        _add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, self.SLIDE_W, 1,
                   fill_color=_c("primary", self.pal))
        tb = _add_textbox(slide, 0.8, 0.12, 11.733, 0.8)
        _add_para(tb.text_frame, title, size=28, bold=True,
                  color=RGBColor(255, 255, 255), align=PP_ALIGN.LEFT)

        cd = CategoryChartData()
        cd.categories = categories
        for s in series_data:
            cd.add_series(s.get("name", ""), s.get("values", []))

        ct_map = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "line": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
        }
        ct = ct_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_shape = slide.shapes.add_chart(
            ct, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2), cd
        )
        chart = chart_shape.chart
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        if hasattr(chart, 'font'):
            chart.font.size = Pt(10)

        if chart_type == "pie":
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.show_percentage = True
            data_labels.show_value = False
            data_labels.font.size = Pt(10)

        return slide


# ── JSON input → PPT pipeline ────────────────────────────
def build_from_json(data):
    style = data.get("style", "business_blue")
    gen = PPTGenerator(
        style=style,
        footer=data.get("footer", ""),
        page_numbers=data.get("page_numbers", True),
    )
    total = len(data.get("slides", []))
    for idx, s in enumerate(data.get("slides", [])):
        typ = s.get("type", "content")
        title = s.get("title", "")
        try:
            if typ == "cover":
                gen.add_cover(
                    title=title,
                    subtitle=s.get("subtitle", ""),
                    author=s.get("author", ""),
                    date=s.get("date", ""),
                    variant=s.get("variant", "centered"),
                )
            elif typ == "section":
                gen.add_section(title=title, subtitle=s.get("subtitle"))
            elif typ == "content":
                gen.add_content(
                    title=title,
                    items=s.get("items", []),
                    columns=s.get("columns", 1),
                )
            elif typ == "two_column":
                gen.add_two_column(
                    title=title,
                    left_items=s.get("left", []),
                    right_items=s.get("right", []),
                )
            elif typ == "table":
                gen.add_table(
                    title=title,
                    headers=s.get("headers", []),
                    rows=s.get("rows", []),
                )
            elif typ == "image":
                gen.add_image_slide(
                    title=title,
                    image_path=s.get("image_path", ""),
                    caption=s.get("caption", ""),
                    overlay=s.get("overlay", False),
                )
            elif typ == "summary":
                gen.add_summary(
                    title=title,
                    points=s.get("points", []),
                    conclusion=s.get("conclusion", ""),
                )
            elif typ == "timeline":
                gen.add_timeline(
                    title=title,
                    milestones=s.get("milestones", []),
                )
            elif typ == "quote":
                gen.add_quote(
                    quote_text=s.get("quote", ""),
                    attribution=s.get("attribution", ""),
                )
            elif typ == "contact":
                gen.add_contact(
                    title=title,
                    contact_info=s.get("info", ""),
                )
            elif typ == "agenda":
                gen.add_agenda(
                    title=title,
                    topics=s.get("topics", []),
                    highlight=s.get("highlight", 0),
                )
            elif typ == "chart":
                gen.add_chart(
                    title=title,
                    chart_type=s.get("chart_type", "bar"),
                    categories=s.get("categories", []),
                    series_data=s.get("series", []),
                )
        except Exception as e:
            print(f"[{idx + 1}/{total}] 跳过 \"{title}\" ({typ}): {e}", file=sys.stderr)
    return gen


def main():
    args = sys.argv[1:]
    if not args or "--help" in args or "-h" in args:
        print("""PPT 生成器 v2.0.0

用法:
  python3 generate_pptx.py <json_file> [--out output.pptx] [--style business_blue]

JSON 结构:
  {
    "style": "business_blue",
    "footer": "机密",
    "page_numbers": true,
    "slides": [
      { "type": "cover", "title": "...", "subtitle": "...", "variant": "centered" },
      { "type": "section", "title": "...", "subtitle": "..." },
      { "type": "agenda", "title": "目录", "topics": ["主题1", "主题2"], "highlight": 0 },
      { "type": "content", "title": "...", "items": ["项1", "项2"], "columns": 1 },
      { "type": "two_column", "title": "...", "left": ["a"], "right": ["b"] },
      { "type": "table", "title": "...", "headers": ["A","B"], "rows": [[1,2]] },
      { "type": "chart", "title": "...", "chart_type": "bar|line|pie",
        "categories": ["Q1","Q2"], "series": [{"name":"S1","values":[10,20]}] },
      { "type": "timeline", "title": "...",
        "milestones": [{"label":"阶段1","desc":"描述"}] },
      { "type": "quote", "title": "...", "quote": "内容", "attribution": "作者" },
      { "type": "image", "title": "...", "image_path": "img.png", "caption": "...", "overlay": false },
      { "type": "summary", "title": "...", "points": ["要点"], "conclusion": "结论" },
      { "type": "contact", "title": "联系我们", "info": "email@example.com" }
    ]
  }

风格: business_blue, academic_white, creative_purple, tech_dark, minimal_gray""")
        return

    json_path = args[0]
    flags = {}
    for i, a in enumerate(args):
        if a == "--out" and i + 1 < len(args):
            flags["out"] = args[i + 1]
        if a == "--style" and i + 1 < len(args):
            flags["style"] = args[i + 1]

    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    if "style" not in data and flags.get("style"):
        data["style"] = flags["style"]

    gen = build_from_json(data)

    out = flags.get("out") or data.get("output", "output.pptx")
    gen.prs.save(out)
    print(f"✅ 已生成 {out}  ({gen._page} 页)")


if __name__ == "__main__":
    try:
        main()
    except FileNotFoundError as e:
        print(f"错误: {e}", file=sys.stderr)
        sys.exit(1)
    except json.JSONDecodeError as e:
        print(f"JSON 解析错误: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"错误: {e}", file=sys.stderr)
        sys.exit(1)
